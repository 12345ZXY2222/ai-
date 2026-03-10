
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation(output_file):
    prs = Presentation()

    # Helper to add a slide with title and content
    def add_slide(title_text, content_text_list):
        slide_layout = prs.slide_layouts[1] # Bullet slide
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = title_text
        
        tf = body_shape.text_frame
        tf.text = content_text_list[0]
        
        for item in content_text_list[1:]:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0

    # Helper for title slide
    def add_title_slide(title_text, subtitle_text):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = title_text
        subtitle.text = subtitle_text

    # 1. Title Slide
    add_title_slide(
        "大语言模型对人类经济社会行为的模拟",
        "基于生成式智能体的计算社会科学研究\n汇报人：[您的名字]\n日期：2025年12月14日"
    )

    # 2. 研究背景与动机
    add_slide(
        "1. 研究背景与动机",
        [
            "社会科学研究的困境：",
            "  - 实验室实验：样本小、环境人工化、成本高。",
            "  - 实证研究：因果推断困难。",
            "  - 传统ABM：微观基础薄弱，缺乏认知深度。",
            "解决方案：生成式智能体 (Generative Agents)",
            "  - 利用LLM作为“硅基被试” (Silicon Subjects)。",
            "  - 具备海量知识、推理能力与自然语言交互能力。",
            "  - 低成本、高可控、大规模并发。"
        ]
    )

    # 3. 实验平台架构
    add_slide(
        "2. 通用社会模拟实验平台",
        [
            "系统架构：",
            "  - 前端：React + 可视化仿真设计器。",
            "  - 后端：FastAPI + 智能体生成器 (Persona Injection)。",
            "核心引擎：",
            "  - 架构师 (Architect)：解析自然语言需求。",
            "  - 开发者 (Developer)：生成可执行步骤 (JSON)。",
            "  - 执行引擎：动态上下文解析 + 代码沙箱。",
            "数据层：",
            "  - 完整的日志记录与状态回溯。"
        ]
    )

    # 4. 实验一：AI道德水平基准测试
    add_slide(
        "3.1 实验一：AI道德水平基准测试",
        [
            "实验范式：",
            "  - 最后通牒博弈、独裁者博弈、公共品博弈、信任博弈。",
            "  - 对比 DeepSeek 与 Qwen 模型。",
            "  - 变量：基准组 (中性) vs 框架组 (道德情境)。",
            "主要发现：",
            "  - 道德异质性：DeepSeek 倾向于规则导向的绝对公平；Qwen 对社会情境敏感，表现出互惠行为。",
            "  - 框架效应：Qwen 在“慈善”框架下捐赠额大幅提升；DeepSeek 对“搭便车”逻辑反应敏感（转向背叛）。"
        ]
    )

    # 5. 实验二：库存管理博弈
    add_slide(
        "3.2 实验二：库存管理博弈",
        [
            "实验设置：",
            "  - 单产品多周期库存管理 (报童模型扩展)。",
            "  - 考察长鞭效应 (Bullwhip Effect) 与动态规划能力。",
            "结果：",
            "  - 基础表现：存在严重的反应滞后与库存积压，利润远低于最优策略。",
            "  - 知识增强效应：引入 CoT (思维链) + 领域知识后，DeepSeek 表现跃升至 99.5% 最优水平。",
            "  - 启示：LLM 需要特定领域的知识注入才能激活其推理能力。"
        ]
    )

    # 6. 实验三：金融市场与泡沫
    add_slide(
        "3.3 实验三：金融市场与泡沫",
        [
            "市场构成：",
            "  - 50个智能体：基本面交易者、技术交易者、社交交易者 (LLM)。",
            "  - 机制：小世界网络传播信息，做市商定价。",
            "现象：",
            "  - 意见领袖：LLM 智能体加速了群体情绪的传播。",
            "  - 泡沫形态：高波动环境下自发崩盘；低波动环境下受冲击驱动。",
            "  - 财富再分配：LLM 在高波动中表现出风险厌恶，但在低波动中容易过度自信被套牢。"
        ]
    )

    # 7. 核心发现：理性悖论
    add_slide(
        "4. 核心发现：理性悖论",
        [
            "CoT 干预的双面性：",
            "  - 宏观层面：CoT + 领域知识 有效平抑了市场泡沫，降低了系统波动率。",
            "  - 微观层面：理性的智能体 (CoT组) 在非理性市场中反而亏损。",
            "解释：",
            "  - “市场保持非理性的时间，可能比你保持偿付能力的时间更长。”",
            "  - 在充满噪声交易者的生态中，个体的绝对理性并不总是最优生存策略。"
        ]
    )

    # 8. 结论与展望
    add_slide(
        "5. 结论与展望",
        [
            "结论：",
            "  - LLM 是有效的社会模拟工具，能够复现经典经济学现象。",
            "  - 模型间存在显著的性格与能力差异 (DeepSeek vs Qwen)。",
            "  - 提示词工程 (Prompt Engineering) 对模型表现至关重要。",
            "展望：",
            "  - 提升大规模多智能体模拟的计算效率。",
            "  - 引入因果推断框架验证模拟结果。",
            "  - 探索更多复杂的社会系统 (如政策制定、舆情传播)。"
        ]
    )

    prs.save(output_file)
    print(f"Presentation saved to {output_file}")

if __name__ == "__main__":
    create_presentation("AI_Social_Simulation_Report.pptx")
